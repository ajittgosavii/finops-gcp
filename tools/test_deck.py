"""The deck must be self-explanatory, and must fit on the slide.

Two failures have already shipped here and both were invisible to `python
tools/build_deck.py` exiting 0: a card row that put its fourth card 3.3 inches
off a 13.33-inch slide, and a bullet block that printed straight through the
caption. "It wrote the file" is not evidence that anyone can read it.

And the deck gets emailed. Whoever opens it will not have the presenter, so
every slide has to carry its own argument in the notes pane.

    DATA_SOURCE=demo PYTHONPATH=services/api pytest tools
"""

from __future__ import annotations

import os
import sys

import pytest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, "tools"))
sys.path.insert(0, os.path.join(ROOT, "packages", "finops-core", "src"))

pytest.importorskip("pptx")
pytest.importorskip("matplotlib")

import build_deck as bd  # noqa: E402
from pptx import Presentation  # noqa: E402

EXPECTED_SLIDES = 23
DIAGRAM_SLIDES = (5, 6, 7, 8)  # hld, end_user_view, lld, cloud_onboarding


@pytest.fixture(scope="module")
def deck(tmp_path_factory):
    out = str(tmp_path_factory.mktemp("deck") / "test.pptx")
    bd.build(out)
    return Presentation(out)


def test_the_deck_has_the_slides_we_think_it_has(deck) -> None:
    assert len(deck.slides) == EXPECTED_SLIDES


def test_every_slide_carries_speaker_notes(deck) -> None:
    """The presenter may not be in the room. A slide with no notes is a slide
    that cannot explain itself to whoever opens the file."""
    thin = []
    for i, slide in enumerate(deck.slides, 1):
        text = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        if len(text) < 200:
            thin.append((i, len(text)))
    assert not thin, f"slides with missing or thin speaker notes: {thin}"


def test_the_notes_list_cannot_drift_out_of_alignment() -> None:
    """`speaker_notes` is one list in build order. A note attached to the wrong
    slide is worse than no note, so the count is asserted at build time too."""
    notes = bd.speaker_notes(bd.gather())
    assert len(notes) == EXPECTED_SLIDES


def test_nothing_falls_off_the_slide(deck) -> None:
    """Every shape sits inside the canvas. This is the assertion that would have
    caught the FOCUS-emitter cards running 3.3in past the right edge."""
    W, H = deck.slide_width, deck.slide_height
    off = []
    for i, slide in enumerate(deck.slides, 1):
        for sh in slide.shapes:
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > W or sh.top + sh.height > H:
                off.append((i, sh.shape_type, round(sh.left / 914400, 2), round((sh.left + sh.width) / 914400, 2)))
    assert not off, f"shapes outside the canvas: {off}"


def test_the_diagram_slides_carry_a_how_to_read_it_rail(deck) -> None:
    for n in DIAGRAM_SLIDES:
        slide = list(deck.slides)[n - 1]
        text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert "HOW TO READ IT" in text, f"slide {n} lost its rail"


def test_the_rail_never_overlaps_the_diagram(deck) -> None:
    """The rail lives in the dead space beside the picture. If a diagram's aspect
    ratio changes, the picture could grow into it -- so check, do not assume."""
    for n in DIAGRAM_SLIDES:
        slide = list(deck.slides)[n - 1]
        pic = next(sh for sh in slide.shapes if sh.shape_type == 13)
        rail = next(sh for sh in slide.shapes
                    if sh.has_text_frame and sh.text_frame.text.strip() == "HOW TO READ IT")
        assert rail.left >= pic.left + pic.width, f"slide {n}: rail overlaps the picture"


def test_every_diagram_slide_names_its_editable_vector_source(deck) -> None:
    for n in DIAGRAM_SLIDES:
        slide = list(deck.slides)[n - 1]
        text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert "docs/diagrams/" in text and ".svg" in text
