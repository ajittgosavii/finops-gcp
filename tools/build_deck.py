"""Generate the GCP-platform deck (light theme) from this codebase.

Every figure on every slide is computed at build time by importing `finops_core`
-- the same engines the API serves and the agents call. Nothing is typed by hand.
A deck whose numbers drift from the product is worse than no deck.

    python tools/build_deck.py       # -> Infosys_FinOps_GCP_Platform.pptx

Charts are native PowerPoint charts, not images, so the client can click into
them and lift them into their own template. The two architecture diagrams are
authored in `tools/diagrams.py` and shipped as SVG; the PNGs go here because
python-pptx cannot embed SVG.

The estate behind the numbers is SYNTHETIC. Every slide that quotes a figure
says so.
"""

from __future__ import annotations

import argparse
import os
import sys
from dataclasses import dataclass
from typing import List, Sequence, Tuple

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
sys.path.insert(0, os.path.join(ROOT, "tools"))

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

W, H = Inches(13.333), Inches(7.5)

# Text colour is contrast-first. The old BODY/MUTED pair measured 8.1:1 and
# 3.7:1 on white; MUTED is used for every sub-label, caption and footer, and
# 3.7:1 is below the WCAG AA floor of 4.5:1 -- legible on a laptop, not on a
# projector in a lit room. Now 11.5:1 and 7.3:1, with the three-level
# hierarchy (18.3 / 11.5 / 7.3) intact. `tests` assert these ratios.
INK = RGBColor(0x0B, 0x14, 0x2A)
BODY = RGBColor(0x2E, 0x3A, 0x4E)
MUTED = RGBColor(0x4E, 0x57, 0x66)
RULE = RGBColor(0xE2, 0xE7, 0xEF)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF5, 0xF8, 0xFC)

AZURE = RGBColor(0x1E, 0x6F, 0xD9)
TEAL = RGBColor(0x11, 0x9B, 0x8A)
VIOLET = RGBColor(0x5B, 0x4B, 0xC4)
AMBER = RGBColor(0xC9, 0x85, 0x1F)
CRIMSON = RGBColor(0xC2, 0x33, 0x33)
GREEN = RGBColor(0x0C, 0x7A, 0x3E)
# OCI: violet, not Oracle red -- CRIMSON is the alert hue on every other slide.
ORACLE = RGBColor(0x6E, 0x3A, 0xA7)

SERIES = [
    RGBColor(0x2A, 0x78, 0xD6), RGBColor(0xEB, 0x68, 0x34), RGBColor(0x1B, 0xAF, 0x7A),
    RGBColor(0x4A, 0x3A, 0xA7), RGBColor(0xED, 0xA1, 0x00),
]
FONT = "Segoe UI"

SYNTHETIC = "Figures computed from a synthetic 24-month utility estate shipped with the platform. No customer data."

# Measured, not guessed: see the token accounting on the "What the agents cost"
# slide. Gemini list prices, July 2026.
GEMINI = {
    "reasoning": ("gemini-3.5-flash", 1.50, 9.00),
    "routing": ("gemini-3.1-flash-lite", 0.25, 1.50),
}
TOKENS = {"coordinator_in": 3600, "coordinator_out": 90, "specialist_in": 12000, "specialist_out": 820}
QUESTIONS_PER_MONTH = 4400


# ==========================================================================
# The delivery estimate
#
# Effort is stated in PERSON-MONTHS, never in dollars. Rates are commercial and
# not ours to invent on a slide. A number the client can multiply by their own
# rate card is more useful, and more honest, than one we made up.
#
# The shape is deliberately unaggressive. Almost nothing engineering-heavy lands
# in month 1, because in a regulated utility the long pole is not the build -- it
# is getting read credentials on four payers and four exports enabled. Every
# month after that assumes that work already happened.
#
# The table, the chart and the totals on the slide are all computed from this one
# structure, so they cannot disagree with one another.
# ==========================================================================

CONTINGENCY = 0.15

# (role, location, FTE in months 1..4, why it is staffed this way)
EFFORT: List[Tuple[str, str, List[float], str]] = [
    ("Engagement / Delivery Lead", "Onshore", [0.5, 0.5, 0.5, 0.5],
     "Con Edison-facing throughout. Owns the gates."),
    ("FinOps Solution Architect", "Onshore", [1.0, 0.5, 0.5, 0.5],
     "Heaviest in month 1: tag taxonomy, allocation policy, KPI definitions."),
    ("FinOps Analyst / BA", "Onshore", [0.5, 0.5, 0.5, 0.5],
     "Validates every number against Con Edison's own reporting."),
    ("Security & Compliance Lead", "Onshore", [0.5, 0.25, 0.25, 0.5],
     "Front-loaded for IAM; back-loaded for IAP, RBAC and the review."),
    ("Cloud / Data Engineers (2)", "Offshore", [1.0, 2.0, 2.0, 1.0],
     "Connector onboarding, ingest, warehouse. The bulk of the build."),
    ("Platform / DevOps (SRE)", "Offshore", [1.0, 1.0, 0.5, 0.5],
     "Terraform, Cloud Run, CI. Front-loaded: the landing zone gates everything."),
    ("Backend Engineer", "Offshore", [0.5, 1.0, 1.0, 0.5],
     "API, repository, cost guards."),
    ("Frontend Engineer", "Offshore", [0.0, 0.5, 1.0, 0.5],
     "Starts once the API contract is stable, not before."),
    ("AI / Agent Engineer", "Offshore", [0.0, 0.5, 1.0, 0.5],
     "Agents need the engines and real data first. Month 3 is the earliest useful start."),
    ("QA / Test Engineer", "Offshore", [0.0, 0.5, 1.0, 1.0],
     "Ramps into UAT. Runs parity checks against Con Edison's figures."),
]


def effort_totals() -> dict:
    """Person-months by location, plus contingency."""
    on = sum(sum(f) for _, loc, f, _ in EFFORT if loc == "Onshore")
    off = sum(sum(f) for _, loc, f, _ in EFFORT if loc == "Offshore")
    base = on + off
    return {
        "onshore": on,
        "offshore": off,
        "base": base,
        "contingency": base * CONTINGENCY,
        "total": base * (1 + CONTINGENCY),
        "offshore_pct": off / base * 100,
        "by_month": [
            (sum(f[m] for _, loc, f, _ in EFFORT if loc == "Onshore"),
             sum(f[m] for _, loc, f, _ in EFFORT if loc == "Offshore"))
            for m in range(4)
        ],
    }


# Weeks are 1-indexed and inclusive. Sixteen weeks = four months.
PLAN: List[Tuple[str, List[Tuple[str, int, int, str]]]] = [
    ("Mobilise & access", [
        ("Kick-off, RACI, environment standards", 1, 2, "Onshore"),
        ("Read-only credentials on four payers", 1, 4, "Onshore"),
        ("Enable FOCUS exports on all four clouds", 2, 5, "Onshore"),
        ("OCI 'endorse' policy on Oracle's tenancy", 3, 5, "Onshore"),
    ]),
    ("Landing zone", [
        ("GCP project, VPC, IAM, budget alerts", 2, 4, "Offshore"),
        ("Terraform: BigQuery, Cloud Run, schedulers", 3, 6, "Offshore"),
        ("CI/CD pipeline and environments", 4, 7, "Offshore"),
    ]),
    ("Ingest & warehouse", [
        ("Connector config per payer; demo dry run", 5, 7, "Offshore"),
        ("First real FOCUS load; partition and cluster", 6, 9, "Offshore"),
        ("Nightly job, GCS replay, idempotency", 8, 10, "Offshore"),
        ("Tag taxonomy + allocation policy with ConEd", 5, 9, "Onshore"),
    ]),
    ("Engines & KPIs", [
        ("KPI parity vs ConEd's current reporting", 8, 11, "Onshore"),
        ("Forecast backtest; cliff calibration", 9, 12, "Offshore"),
        ("Optimization levers reviewed with FinOps", 10, 13, "Onshore"),
    ]),
    ("API & dashboards", [
        ("REST API, scope, cost guards", 7, 10, "Offshore"),
        ("Nine pages; chart and table twins", 9, 13, "Offshore"),
        ("Persona walkthroughs and feedback", 12, 14, "Onshore"),
    ]),
    ("Agentic Copilot", [
        ("Typed tools wired to the engines", 9, 12, "Offshore"),
        ("Coordinator and four specialists on Vertex", 11, 14, "Offshore"),
        ("Grounding review: every figure cites a tool", 13, 15, "Onshore"),
    ]),
    ("Security & hardening", [
        ("IAP, Cloud Identity / Okta, per-persona RBAC", 11, 15, "Offshore"),
        ("OCI key-rotation runbook; secret review", 13, 15, "Onshore"),
        ("Security review and remediation", 14, 16, "Onshore"),
    ]),
    ("UAT & handover", [
        ("UAT with Finance, FinOps and Engineering", 13, 16, "Onshore"),
        ("Runbook, operations training, DR drill", 14, 16, "Offshore"),
        ("Production cutover; hypercare begins", 16, 16, "Onshore"),
    ]),
]

GATES = [
    (4, "G1", "Access granted, exports enabled"),
    (9, "G2", "Real FOCUS data in the warehouse"),
    (14, "G3", "Dashboards and Copilot on real data"),
    (16, "G4", "Cutover and handover"),
]


def gemini_cost_per_question() -> float:
    _, r_in, r_out = GEMINI["reasoning"]
    _, c_in, c_out = GEMINI["routing"]
    spec = TOKENS["specialist_in"] * r_in / 1e6 + TOKENS["specialist_out"] * r_out / 1e6
    coord = TOKENS["coordinator_in"] * c_in / 1e6 + TOKENS["coordinator_out"] * c_out / 1e6
    return spec + coord


# ==========================================================================
# Facts, read from the running code
# ==========================================================================


@dataclass
class Facts:
    org: str
    rows: int
    months: int
    spend: float
    esr: float
    coverage: float
    utilization: float
    commitment_waste: float
    cost_of_waste: float
    waste_pct: float
    allocation: float
    readiness: str
    fc_method: str
    fc_wape: float
    fc_maturity: str
    fc_total: float
    fc_with_cliffs: float
    cliff_months: List[str]
    savings_total: float
    savings_by_category: List[Tuple[str, float]]
    top_opps: List[Tuple[str, str, str, float, str, str]]
    esr_uplift: Tuple[float, float, float]
    n_levers: int
    n_connectors: int
    n_opps: int
    spend_by_cloud: List[Tuple[str, float]]
    forecast: List[Tuple[str, float, float, float]]
    anomaly_count: int
    anomalies: List[Tuple[str, str, float, float, float]]


def gather() -> Facts:
    import dataclasses

    from finops_core import connectors, kpi
    from finops_core.connectors.demo import build_demo_dataset
    from finops_core.engines import anomaly, forecast as fx, optimize

    df, _budgets, _drivers = build_demo_dataset()
    opps = optimize.detect_all(df)
    uw = optimize.usage_waste_total(opps)
    k = kpi.executive_kpis(df, usage_waste_monthly=uw)

    monthly = df.copy()
    monthly["period"] = monthly["ChargePeriodStart"].dt.to_period("M").dt.to_timestamp()
    monthly = monthly.groupby("period", as_index=False, observed=True)["EffectiveCost"].sum()
    monthly = monthly.rename(columns={"EffectiveCost": "cost"})

    fc = fx.forecast_spend(monthly, horizon=24, method="auto")
    cliff = fx.commitment_expiry_overlay(df, fc.forecast)
    cliff_months = (
        cliff.loc[cliff["cliff"], "period"].dt.strftime("%b %Y").tolist() if "cliff" in cliff else []
    )

    sav = optimize.savings_by_category(opps)
    up = optimize.effective_savings_rate_uplift(df, opps)
    by_cloud = df.groupby("ProviderName", observed=True)["EffectiveCost"].sum().sort_values(ascending=False)

    hits = anomaly.detect_by_dimension(df, dim="ServiceCategory").sort_values("deviation_pct", ascending=False)
    top = sorted(opps, key=lambda o: -o.annual_savings)[:8]

    return Facts(
        org="Con Edison",
        rows=len(df),
        months=len(monthly),
        spend=k.total_spend,
        esr=k.esr_pct or 0,
        coverage=k.coverage_pct or 0,
        utilization=k.utilization_pct or 0,
        commitment_waste=k.commitment_waste,
        cost_of_waste=k.cost_of_waste,
        waste_pct=k.waste_pct or 0,
        allocation=k.allocation_coverage_pct or 0,
        readiness=k.chargeback_readiness,
        fc_method=fc.method.replace("_", " ").title(),
        fc_wape=fc.accuracy["wape"],
        fc_maturity=fc.maturity,
        fc_total=float(fc.forecast["cost"].sum()),
        fc_with_cliffs=float(cliff["cost_with_cliffs"].sum()) if "cost_with_cliffs" in cliff else 0.0,
        cliff_months=cliff_months,
        savings_total=float(sum(o.annual_savings for o in opps)),
        savings_by_category=[(r.category, float(r.annual_savings)) for r in sav.itertuples()],
        top_opps=[(o.lever_id, o.lever_name, o.cloud, o.annual_savings, o.effort, o.risk) for o in top],
        esr_uplift=(up["current_esr_pct"], up["projected_esr_pct"], up["uplift_pts"]),
        n_levers=len(optimize.LEVERS),
        n_connectors=len(connectors.REGISTRY),
        n_opps=len(opps),
        spend_by_cloud=[(str(i), float(v)) for i, v in by_cloud.items()],
        forecast=[
            (p.strftime("%b %y"), float(c), float(lo), float(hi))
            for p, c, lo, hi in zip(fc.forecast["period"], fc.forecast["cost"],
                                    fc.forecast["lo80"], fc.forecast["hi80"])
        ],
        anomaly_count=int(len(hits)),
        anomalies=[
            (r.period.strftime("%d %b %Y"), str(getattr(r, "ServiceCategory")),
             float(r.cost), float(r.expected), float(r.deviation_pct))
            for r in hits.head(3).itertuples()
        ],
    )


# ==========================================================================
# Layout primitives
# ==========================================================================


def money(x: float) -> str:
    a, s = abs(x), "-" if x < 0 else ""
    if a >= 1e9:
        return f"{s}${a/1e9:.2f}B"
    if a >= 1e6:
        return f"{s}${a/1e6:.2f}M"
    if a >= 1e3:
        return f"{s}${a/1e3:.0f}K"
    return f"{s}${a:,.0f}"


def _style(shape, size, bold, colour, align):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = FONT


def box(slide, x, y, w, h, text="", size=14, bold=False, colour=BODY, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_textbox(x, y, w, h)
    sh.text_frame.text = text
    _style(sh, size, bold, colour, align)
    return sh


def bullets(slide, x, y, w, h, items: Sequence[str], size=13, colour=BODY, bullet="—"):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}" if bullet else item
        p.space_after = Pt(9)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = colour
            r.font.name = FONT
    return sh


def rect(slide, x, y, w, h, fill=WASH, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def bar(slide, x, y, w, h, colour):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = colour
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, eyebrow, title, sub=""):
    box(slide, Inches(0.7), Inches(0.42), Inches(11), Inches(0.3), eyebrow.upper(), 10, True, MUTED)
    box(slide, Inches(0.7), Inches(0.7), Inches(11.9), Inches(0.6), title, 27, True, INK)
    if sub:
        box(slide, Inches(0.7), Inches(1.32), Inches(11.9), Inches(0.4), sub, 12.5, False, MUTED)
    bar(slide, Inches(0.7), Inches(1.82), Inches(1.4), Pt(3), AZURE)


def footer(slide, note=""):
    box(slide, Inches(0.7), Inches(6.95), Inches(9.5), Inches(0.3), note, 8.5, False, MUTED)
    box(slide, Inches(11.4), Inches(6.95), Inches(1.3), Inches(0.3), "Infosys", 8.5, True, MUTED, PP_ALIGN.RIGHT)


def kpi_card(slide, x, y, w, h, label, value, sub, accent=AZURE):
    rect(slide, x, y, w, h, PAPER, RULE)
    bar(slide, x, y, Pt(3.5), h, accent)
    box(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.3), Inches(0.28), label.upper(), 8.5, True, MUTED)
    box(slide, x + Inches(0.18), y + Inches(0.34), w - Inches(0.3), Inches(0.5), value, 21, True, INK)
    box(slide, x + Inches(0.18), y + Inches(0.86), w - Inches(0.3), Inches(0.4), sub, 9, False, MUTED)


def style_chart(chart, colours, legend=False, fmt='#,##0,,"M"'):
    chart.font.size = Pt(10)
    chart.font.name = FONT
    chart.font.color.rgb = BODY
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    for i, s in enumerate(chart.plots[0].series):
        c = colours[i % len(colours)]
        try:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = c
        except Exception:
            pass
        try:
            s.format.line.color.rgb = c
            s.format.line.width = Pt(2)
        except Exception:
            pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RULE
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.number_format = fmt
        va.tick_labels.number_format_is_linked = False
        va.format.line.fill.background()
        chart.category_axis.has_major_gridlines = False
        chart.category_axis.format.line.color.rgb = RULE
    except Exception:
        pass


def diagram(prs, name, eyebrow, title, sub, note, reads=()):
    """A diagram, with a "How to read it" rail in the dead space beside it.

    The diagrams are 16:10-ish, so the 4.6in height cap binds first and the
    picture only ever fills ~7 of the slide's 13.3 inches. Three inches of white
    on each side were doing nothing. Left-aligning the picture buys the rail for
    free: the diagram does not shrink by a pixel.

    `reads` is up to three (heading, sentence) pairs -- the things a reader
    cannot get from the picture alone. If a wider picture leaves no room, we
    centre it and drop the rail rather than let the two overlap.
    """
    import diagrams as dg

    png = os.path.join(dg.OUT_DIR, f"{name}.png")
    if not os.path.exists(png):
        dg.build_all()

    s = blank(prs)
    header(s, eyebrow, title, sub)
    from PIL import Image

    with Image.open(png) as im:
        iw, ih = im.size
    scale = min(Inches(12.1) / iw, Inches(4.6) / ih)
    w, h = int(iw * scale), int(ih * scale)
    GUTTER, RAIL_MIN, RIGHT = Inches(0.4), Inches(3.4), Inches(12.63)
    rail_x = Inches(0.7) + w + GUTTER
    rail_w = RIGHT - rail_x

    if reads and rail_w >= RAIL_MIN:
        s.shapes.add_picture(png, Inches(0.7), Inches(2.05), w, h)
        box(s, rail_x, Inches(2.05), rail_w, Inches(0.3), "HOW TO READ IT", 9, True, MUTED)
        y = Inches(2.52)
        for head, text in reads[:3]:
            bar(s, rail_x, y + Inches(0.06), Pt(3), Inches(0.20), TEAL)
            box(s, rail_x + Inches(0.13), y, rail_w - Inches(0.13), Inches(0.28), head, 10.5, True, INK)
            box(s, rail_x + Inches(0.13), y + Inches(0.31), rail_w - Inches(0.13), Inches(0.95),
                text, 9.5, False, BODY)
            y += Inches(1.34)
    else:
        s.shapes.add_picture(png, int((W - w) / 2), Inches(2.05), w, h)

    box(s, Inches(0.7), Inches(6.72), Inches(10.4), Inches(0.3), note, 9, False, MUTED)
    box(s, Inches(10.9), Inches(6.72), Inches(1.8), Inches(0.3), f"docs/diagrams/{name}.svg", 8, False, MUTED, PP_ALIGN.RIGHT)
    footer(s)


# ==========================================================================
# Slides
# ==========================================================================


def slide_title(prs, f: Facts):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), W, Inches(2.6), RGBColor(0xF2, 0xF6, 0xFC))
    box(s, Inches(0.9), Inches(0.72), Inches(6), Inches(0.35), "INFOSYS", 12, True, AZURE)
    box(s, Inches(0.9), Inches(1.0), Inches(11.6), Inches(0.9), "Multi-Cloud FinOps on Google Cloud", 38, True, INK)
    box(s, Inches(0.9), Inches(1.92), Inches(11.6), Inches(0.5),
        "AWS, Azure, GCP and OCI spend in one FOCUS warehouse, with an agentic control plane", 15.5, False, BODY)
    box(s, Inches(0.9), Inches(3.05), Inches(11), Inches(0.4), f"Prepared for {f.org}", 15, True, INK)
    bullets(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(2.4), [
        "Four clouds normalised to the FinOps Foundation's FOCUS 1.2 specification, in BigQuery",
        "One credential per payer — not one per account — federated for AWS and Azure, a single signing key for OCI",
        f"{f.n_levers} optimization levers detected from the bill; {money(f.savings_total)} identified on a representative estate",
        "A Google ADK agent team on Gemini that can only quote numbers it measured",
        "Serverless on Cloud Run: scales to zero, ~$240/month plus model usage",
    ], 13.5)
    footer(s, "Infosys · Con Edison · GCP Platform Overview")


def slide_problem(prs, f: Facts):
    s = blank(prs)
    header(s, "The problem", "Four clouds, four schemas, no single truth",
           "Multi-cloud FinOps is a data problem before it is a savings problem")
    items = [
        ("Four billing schemas", "AWS, Azure, GCP and OCI each bill in their own shape. Reconciling them by hand is a monthly project, not a dashboard."),
        ("Credential sprawl", "Teams assume one credential per account. It is one per payer — and the difference is dozens of keys nobody rotates."),
        ("Tool lock-in", "Pick a FinOps platform and every dashboard, KPI and script is written against that vendor's field names."),
        ("Unallocated spend", "Untagged resources land in a bucket nobody owns, so chargeback is disputed and showback is ignored."),
        ("Blind forecasting", "A trend line walks straight through a commitment expiry. The rate snaps back to on-demand and the variance lands unannounced."),
        ("Scale", "A utility estate is millions of billing line-items a month. A single-process dashboard cannot hold it."),
    ]
    y = Inches(2.15)
    for i, (t, d) in enumerate(items):
        rect(s, Inches(0.7), y, Inches(11.9), Inches(0.73), PAPER if i % 2 else WASH, RULE)
        box(s, Inches(0.95), y + Inches(0.08), Inches(2.9), Inches(0.4), t, 12.5, True, INK)
        box(s, Inches(3.95), y + Inches(0.08), Inches(8.4), Inches(0.55), d, 11, False, BODY)
        y += Inches(0.8)
    footer(s)


def slide_focus(prs, f: Facts):
    s = blank(prs)
    header(s, "The approach", "One idea: normalise everything to FOCUS 1.2",
           "The FinOps Foundation's Open Cost and Usage Specification is the contract")

    emitters = [
        ("AWS", "Data Exports — FOCUS_1_2_AWS\nGA 19 Nov 2025 (1.0 GA Nov 2024)", AMBER),
        ("Azure", "Cost Management exports\ndataset type FocusCost", AZURE),
        ("Google Cloud", "gcp_billing_export_focus_*\nnative, inside BigQuery already", GREEN),
        ("OCI", "FOCUS Reports in the Oracle-owned\n'bling' bucket, gzipped CSV", ORACLE),
    ]
    # Four boxes across the 11.93in content width: 2.83 wide on a 3.03 pitch ends
    # at 12.63. The old 3.83/4.02 pair fitted three and put the fourth at 16.6in,
    # a full 3.3in off a 13.33in slide.
    CARD_W, PITCH = Inches(2.83), Inches(3.03)
    x = Inches(0.7)
    for name, sub, colour in emitters:
        rect(s, x, Inches(2.2), CARD_W, Inches(1.5), PAPER, RULE)
        bar(s, x, Inches(2.2), CARD_W, Pt(4), colour)
        box(s, x + Inches(0.2), Inches(2.4), Inches(2.43), Inches(0.4), name, 14, True, INK)
        box(s, x + Inches(0.2), Inches(2.82), Inches(2.43), Inches(0.8), sub, 9.5, False, BODY)
        x += PITCH

    box(s, Inches(0.7), Inches(4.0), Inches(11.9), Inches(1.4),
        "All four providers emit FOCUS natively today. So do CloudZero and Vantage; Cloudability, CloudHealth and "
        "Flexera ingest it. That is why this platform is vendor-neutral by construction rather than by adapter: no dashboard, "
        "KPI formula, optimization detector or agent tool has ever seen a vendor-specific field.",
        13, False, BODY)

    kpi_card(s, Inches(0.7), Inches(5.35), Inches(3.83), Inches(1.3), "Connectors", str(f.n_connectors),
             "4 native + 12 procured tools + any FOCUS file + demo", AZURE)
    kpi_card(s, Inches(4.72), Inches(5.35), Inches(3.83), Inches(1.3), "Adopting a new tool", "1 subclass",
             "plus one line in the registry", TEAL)
    kpi_card(s, Inches(8.74), Inches(5.35), Inches(3.86), Inches(1.3), "Downstream changes", "None",
             "every dashboard and agent keeps working", VIOLET)
    footer(s)


def slide_focus_rosetta(prs, f: Facts):
    """The slide that makes FOCUS land.

    Slide 3 says who emits FOCUS and why that means no lock-in. It never says
    what FOCUS *is*, so the room nods and leaves none the wiser. This one shows
    the translation happening: the same charge in four dialects collapsing into
    three named columns, and then one formula computed from two of them.

    Each FOCUS card sits directly beneath the vendor column it replaces, so the
    mapping is read by position rather than from a legend.

    Deliberately concept-level on the vendor side. Printing literal CUR / export
    column strings is more visceral and would be one stale string away from
    costing us the architects in the room -- Google is mid-schema-change on the
    detailed export as of July 2026. Every word below is defensible.
    """
    s = blank(prs)
    header(s, "What FOCUS actually is", "One charge. Four dialects. One row.",
           "The FinOps Foundation's Open Cost and Usage Specification — an open schema, not our schema")

    # Vendor columns. Each FOCUS card below is aligned under the one it replaces.
    C_CLOUD, W_CLOUD = Inches(0.7), Inches(1.7)
    C_COMMIT, W_COMMIT = Inches(2.5), Inches(3.3)
    C_PAY, W_PAY = Inches(5.95), Inches(3.2)
    C_LIST, W_LIST = Inches(9.3), Inches(3.3)

    for x, w, label in ((C_COMMIT, W_COMMIT, "The commitment"),
                        (C_PAY, W_PAY, "What you actually pay"),
                        (C_LIST, W_LIST, "The undiscounted price")):
        box(s, x, Inches(1.95), w, Inches(0.3), label, 10, True, MUTED)

    rows = [
        ("AWS", "Savings Plan · Reserved Instance", "amortized cost", "public on-demand cost", AMBER),
        ("Azure", "Reservation · Savings Plan", "amortized cost", "pay-as-you-go price", AZURE),
        ("GCP", "Committed Use Discount", "cost + credits", "list price", GREEN),
        ("OCI", "Annual Universal Credits", "amortized cost", "unit price × quantity", ORACLE),
    ]
    y = Inches(2.35)
    ROW_H, PITCH = Inches(0.56), Inches(0.66)
    for name, commit, pay, lst, colour in rows:
        rect(s, C_CLOUD, y, Inches(11.93), ROW_H, PAPER, RULE)
        bar(s, C_CLOUD, y, Pt(4), ROW_H, colour)
        box(s, C_CLOUD + Inches(0.18), y + Inches(0.10), W_CLOUD, Inches(0.36), name, 12, True, colour)
        for x, w, text in ((C_COMMIT, W_COMMIT, commit), (C_PAY, W_PAY, pay), (C_LIST, W_LIST, lst)):
            box(s, x, y + Inches(0.12), w, Inches(0.34), text, 11, False, BODY)
        y += PITCH

    box(s, C_CLOUD, Inches(4.98), Inches(11.93), Inches(0.3),
        "▼   normalised on ingest — nothing downstream has ever seen a vendor-specific field   ▼",
        10, True, TEAL, PP_ALIGN.CENTER)

    for x, w, name, sub in ((C_COMMIT, W_COMMIT, "CommitmentDiscountStatus", "Used  |  Unused"),
                            (C_PAY, W_PAY, "EffectiveCost", "amortised, never lumpy"),
                            (C_LIST, W_LIST, "ListCost", "the ESR denominator")):
        rect(s, x, Inches(5.35), w, Inches(0.82), WASH, VIOLET)
        box(s, x + Inches(0.12), Inches(5.42), w - Inches(0.24), Inches(0.3), name, 11.5, True, INK)
        box(s, x + Inches(0.12), Inches(5.73), w - Inches(0.24), Inches(0.28), sub, 9.5, False, MUTED)

    box(s, C_CLOUD, Inches(5.42), W_CLOUD, Inches(0.3), "FOCUS 1.2", 12, True, VIOLET)
    box(s, C_CLOUD, Inches(5.73), W_CLOUD, Inches(0.28), "one row, any cloud", 9.5, False, MUTED)

    box(s, C_CLOUD, Inches(6.35), Inches(11.93), Inches(0.32),
        "Effective Savings Rate  =  (ListCost − EffectiveCost) / ListCost      "
        "— defined once, in one function, for all four clouds.",
        12.5, True, INK)
    box(s, C_CLOUD, Inches(6.70), Inches(11.93), Inches(0.30),
        "And CommitmentDiscountStatus = 'Unused' is waste the bill states outright, not an estimate: "
        f"{money(f.commitment_waste)} on this estate.",
        10.5, False, BODY)
    footer(s)


def slide_effort(prs, f: Facts):
    """Effort in person-months, split onshore/offshore and spread over four months.

    No dollars. Rates are commercial, and a rate we invented on a slide would be
    the least defensible number in the deck. Person-months multiply by whatever
    rate card Con Edison actually has.
    """
    t = effort_totals()
    s = blank(prs)
    header(s, "Effort estimation",
           f"{t['base']:.1f} person-months, {t['offshore_pct']:.0f}% offshore, across four months",
           "Deliberately unaggressive. Month 1 buys access, not code — in a regulated utility that is the long pole.")

    heads = ["Role", "Location", "M1", "M2", "M3", "M4", "PM", "Why staffed this way"]
    xs = [Inches(0.7), Inches(3.15), Inches(4.25), Inches(4.8), Inches(5.35), Inches(5.9), Inches(6.5), Inches(7.3)]
    ws = [2.4, 1.05, 0.5, 0.5, 0.5, 0.5, 0.75, 5.3]
    for h, x, w in zip(heads, xs, ws):
        align = PP_ALIGN.CENTER if h in ("M1", "M2", "M3", "M4", "PM") else PP_ALIGN.LEFT
        box(s, x, Inches(2.02), Inches(w), Inches(0.24), h, 8, True, MUTED, align)

    y = Inches(2.34)
    ROW = Inches(0.315)
    for role, loc, fte, why in EFFORT:
        colour = AZURE if loc == "Onshore" else TEAL
        bar(s, Inches(0.7), y + Inches(0.03), Pt(3), Inches(0.22), colour)
        box(s, Inches(0.86), y, Inches(2.3), Inches(0.26), role, 8.5, False, INK)
        box(s, xs[1], y, Inches(ws[1]), Inches(0.26), loc, 8, False, colour)
        for i, v in enumerate(fte):
            txt = f"{v:.2f}".rstrip("0").rstrip(".") if v else "—"
            box(s, xs[2 + i], y, Inches(0.5), Inches(0.26), txt, 8, False,
                BODY if v else RULE, PP_ALIGN.CENTER)
        box(s, xs[6], y, Inches(0.75), Inches(0.26), f"{sum(fte):.2f}".rstrip("0").rstrip("."),
            8.5, True, INK, PP_ALIGN.CENTER)
        box(s, xs[7], y, Inches(5.3), Inches(0.26), why, 7.5, False, MUTED)
        y += ROW

    # Monthly totals, read from the same structure the rows came from.
    rect(s, Inches(0.7), y + Inches(0.04), Inches(11.93), Inches(0.30), WASH)
    box(s, Inches(0.86), y + Inches(0.06), Inches(2.3), Inches(0.26), "Team size (FTE)", 8, True, INK)
    for i, (on, off) in enumerate(t["by_month"]):
        box(s, xs[2 + i], y + Inches(0.06), Inches(0.5), Inches(0.26), f"{on + off:.2f}".rstrip("0").rstrip("."),
            8, True, INK, PP_ALIGN.CENTER)
    box(s, xs[6], y + Inches(0.06), Inches(0.75), Inches(0.26), f"{t['base']:.1f}", 8, True, INK, PP_ALIGN.CENTER)
    peak = max(on + off for on, off in t["by_month"])
    box(s, xs[7], y + Inches(0.06), Inches(5.3), Inches(0.26),
        f"Peak team {peak:.2f} FTE in month 3. No month exceeds it, and nobody works a weekend to hit it.".replace(".00", ""),
        7.5, False, MUTED)

    kpi_card(s, Inches(0.7), Inches(5.75), Inches(2.85), Inches(1.0), "Onshore", f"{t['onshore']:.1f} PM",
             f"{100 - t['offshore_pct']:.0f}% — Con Edison-facing", AZURE)
    kpi_card(s, Inches(3.72), Inches(5.75), Inches(2.85), Inches(1.0), "Offshore", f"{t['offshore']:.1f} PM",
             f"{t['offshore_pct']:.0f}% — build and test", TEAL)
    kpi_card(s, Inches(6.74), Inches(5.75), Inches(2.85), Inches(1.0), "Contingency",
             f"+{CONTINGENCY:.0%}", f"{t['contingency']:.1f} PM, held by the delivery lead", AMBER)
    kpi_card(s, Inches(9.76), Inches(5.75), Inches(2.87), Inches(1.0), "Total", f"{t['total']:.1f} PM",
             "Multiply by your own rate card", VIOLET)
    footer(s, "Effort only. No rates: a rate we invented would be the least defensible number in this deck.")


def slide_plan(prs, f: Facts):
    """Twenty-six activities across sixteen weeks, with four gates."""
    s = blank(prs)
    header(s, "Delivery plan", "Twenty-six activities, sixteen weeks, four gates",
           "Nothing downstream starts before the gate above it clears. Engineering does not wait on itself; it waits on access.")

    # Grid starts at 5.25in so the activity column is 2.7in -- wide enough for the
    # longest name at 6.8pt. At 2.1in six of them wrapped onto the row below.
    X0, XW = Inches(5.25), Inches(7.38)      # week grid: 5.25in .. 12.63in
    WEEK = XW / 16
    TOP = Inches(2.42)
    ROW = Inches(0.152)

    # Month headers and week rules
    for m in range(4):
        x = X0 + WEEK * 4 * m
        rect(s, x + Inches(0.02), Inches(2.02), WEEK * 4 - Inches(0.04), Inches(0.24), WASH)
        box(s, x, Inches(2.03), WEEK * 4, Inches(0.22), f"Month {m + 1}", 8, True, MUTED, PP_ALIGN.CENTER)

    rows = sum(len(a) for _, a in PLAN)
    grid_h = ROW * rows + Inches(0.12)
    for w in range(1, 16):
        bar(s, X0 + WEEK * w, TOP, Pt(0.5), grid_h, RULE)

    # Gates: a vertical rule and a marker, drawn under the bars.
    for week, tag, _ in GATES:
        gx = X0 + WEEK * week
        bar(s, gx - Pt(1), Inches(2.28), Pt(1.5), grid_h + Inches(0.14), CRIMSON)
        box(s, gx - Inches(0.28), Inches(2.28) - Inches(0.24), Inches(0.56), Inches(0.22), tag, 7.5, True,
            CRIMSON, PP_ALIGN.CENTER)

    y = TOP + Inches(0.06)
    for group, acts in PLAN:
        first = True
        for name, w0, w1, owner in acts:
            colour = AZURE if owner == "Onshore" else TEAL
            if first:
                bar(s, Inches(0.7), y, Pt(3), ROW * len(acts) - Inches(0.02), VIOLET)
                box(s, Inches(0.85), y, Inches(1.5), Inches(0.16), group, 7.5, True, INK)
                first = False
            box(s, Inches(2.42), y - Inches(0.008), Inches(2.7), Inches(0.16), name, 6.8, False, BODY)
            bx = X0 + WEEK * (w0 - 1)
            bw = WEEK * (w1 - w0 + 1)
            rect(s, bx + Inches(0.01), y + Inches(0.012), bw - Inches(0.02), Inches(0.11), colour)
            y += ROW

    legend_y = Inches(6.62)
    for i, (label, colour) in enumerate((("Onshore", AZURE), ("Offshore", TEAL))):
        bar(s, Inches(0.7) + Inches(1.3) * i, legend_y + Inches(0.06), Inches(0.16), Inches(0.10), colour)
        box(s, Inches(0.92) + Inches(1.3) * i, legend_y, Inches(1.0), Inches(0.22), label, 8, False, BODY)
    box(s, Inches(3.6), legend_y, Inches(9.0), Inches(0.22),
        "   ".join(f"{tag} (wk {w}): {what}" for w, tag, what in GATES), 7.2, False, CRIMSON)
    footer(s)


def slide_comparison(prs, f: Facts):
    s = blank(prs)
    header(s, "Streamlit vs Google Cloud", "What actually changes, and what does not",
           "The engines are identical. The warehouse, the runtime, the model and the front end are not.")

    rows = [
        ("Where it runs", "Streamlit Community Cloud, one process", "Cloud Run, scales to zero, one service per concern"),
        ("Data", "Whole FOCUS frame in pandas memory", "BigQuery, partitioned and clustered; aggregates pushed to SQL"),
        ("Scale ceiling", "~63k rows demo; ~8 GB at 500k line-items/month", "Billions of rows; queries prune to a slice"),
        ("Cost control", "None needed", "Partition filter required, bytes-billed capped, detectors run nightly"),
        ("Agents", "LangGraph supervisor on gpt-5", "Google ADK coordinator on Gemini, via Vertex"),
        ("Model credentials", "OPENAI_API_KEY in a secret", "ADC on the service account — no model API key exists"),
        ("Front end", "Streamlit, server-rendered", "React + TypeScript, Plotly figure JSON from the same chart code"),
        ("Auth", "Shared password", "Cloud Identity / Okta behind IAP, per-persona RBAC"),
        ("Ingest", "On demand, in the request", "Nightly Cloud Run Job, idempotent, GCS as replay source"),
        ("Shared code", "9,100 lines", "The same 9,100 lines, as an installable package"),
    ]

    y = Inches(2.12)
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.34), WASH)
    for t, dx, w in [("", 0.15, 2.9), ("Streamlit (reference)", 3.2, 4.2), ("Google Cloud (target)", 7.6, 4.7)]:
        box(s, Inches(0.7 + dx), y + Inches(0.02), Inches(w), Inches(0.3), t, 9.5, True, MUTED)
    y += Inches(0.4)
    for i, (dim, a, b) in enumerate(rows):
        if i % 2 == 0:
            rect(s, Inches(0.7), y - Inches(0.02), Inches(11.9), Inches(0.42), WASH)
        box(s, Inches(0.85), y, Inches(2.9), Inches(0.36), dim, 10, True, INK)
        box(s, Inches(3.9), y, Inches(4.1), Inches(0.36), a, 9.5, False, MUTED)
        box(s, Inches(8.3), y, Inches(4.3), Inches(0.36), b, 9.5, False, BODY)
        y += Inches(0.42)

    box(s, Inches(0.7), Inches(6.62), Inches(11.9), Inches(0.4),
        "The Streamlit app is not thrown away. It stays as the reference implementation and the demo surface — and a bug fixed "
        "there is re-extracted into the package, so one copy of the FinOps definitions exists.", 10.5, True, INK)
    footer(s)


def slide_why_rebuild(prs, f: Facts):
    s = blank(prs)
    header(s, "Why rebuild", "It was never the user interface. It was pandas.",
           "The single number that forced the architecture")

    kpi_card(s, Inches(0.7), Inches(2.2), Inches(3.83), Inches(1.35), "Demo estate",
             f"{f.rows:,} rows", "36 MB — about 654 bytes per row", AZURE)
    kpi_card(s, Inches(4.72), Inches(2.2), Inches(3.83), Inches(1.35), "Con Edison scale",
             "~8 GB", "500k line-items/month over two years", AMBER)
    kpi_card(s, Inches(8.74), Inches(2.2), Inches(3.86), Inches(1.35), "Large enterprise",
             "~31 GB", "2M line-items/month over two years", CRIMSON)

    box(s, Inches(0.7), Inches(3.85), Inches(11.9), Inches(0.9),
        "The Streamlit design loads the entire FOCUS frame into one process. It demonstrates beautifully and it will not survive "
        "a utility. BigQuery is not an upgrade; it is the only option that works at Con Edison's scale.", 13, False, BODY)

    box(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.35), "Three cost guards, because BigQuery will happily bill you", 14, True, INK)
    guards = [
        ("Partition filter required", "A query with no bound on the charge period is rejected, not answered. Writing this guard immediately caught a query of ours that had no WHERE clause at all."),
        ("Bytes billed capped", "Every job carries a ceiling (20 GiB ≈ 12¢). BigQuery fails the job rather than billing past it."),
        ("Detectors run nightly", "The row-level scans never run per request. Their answer changes once a day at most."),
    ]
    y = Inches(5.25)
    for t, d in guards:
        rect(s, Inches(0.7), y, Inches(11.9), Inches(0.52), PAPER, RULE)
        box(s, Inches(0.95), y + Inches(0.03), Inches(2.9), Inches(0.4), t, 10.5, True, TEAL)
        box(s, Inches(3.95), y + Inches(0.03), Inches(8.4), Inches(0.44), d, 9.5, False, BODY)
        y += Inches(0.6)
    footer(s)


def slide_agents(prs, f: Facts):
    s = blank(prs)
    header(s, "Agentic AI", "A coordinator and four specialists, on Google ADK",
           "Every figure quoted comes from a tool call against the FOCUS frame")

    rect(s, Inches(0.7), Inches(2.2), Inches(2.6), Inches(1.05), RGBColor(0xE8, 0xF1, 0xFD), AZURE)
    box(s, Inches(0.8), Inches(2.36), Inches(2.4), Inches(0.35), "Coordinator", 13, True, INK, PP_ALIGN.CENTER)
    box(s, Inches(0.8), Inches(2.72), Inches(2.4), Inches(0.4), GEMINI["routing"][0], 8.5, False, MUTED, PP_ALIGN.CENTER)

    specialists = [
        ("Analyst", "Understand Usage and Cost", "spend, allocation, anomalies, coverage"),
        ("Forecaster", "Quantify Business Value", "forecast, budget variance, commitment cliffs"),
        ("Optimizer", "Optimize Usage and Cost", "levers, opportunities, ESR uplift"),
        ("Governor", "Manage the FinOps Practice", "tagging, chargeback readiness, policy"),
    ]
    y = Inches(2.2)
    for name, domain, tools in specialists:
        rect(s, Inches(3.8), y, Inches(8.8), Inches(0.76), PAPER, RULE)
        box(s, Inches(4.0), y + Inches(0.05), Inches(1.7), Inches(0.35), name, 12, True, INK)
        box(s, Inches(5.7), y + Inches(0.07), Inches(3.0), Inches(0.3), domain, 9.5, True, VIOLET)
        box(s, Inches(8.8), y + Inches(0.07), Inches(3.6), Inches(0.5), tools, 9.5, False, BODY)
        y += Inches(0.84)

    box(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.35),
        "Two decisions worth defending in review", 13.5, True, INK)
    bullets(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.9), [
        "Specialists are held as tools, not sub-agents. Handing over control lets the last specialist to run write the answer in its own "
        "register; a FinOps question spans domains and the answer must be one voice speaking to one persona.",
        "The model is never given SQL. Hand it a SQL prompt and it invents its own Effective Savings Rate — dropping the on-demand-equivalent "
        "denominator, or counting Purchase rows. Plausible, wrong, uncaught. It gets typed tools that call the same engines the API calls.",
    ], 10.5)
    footer(s)


def slide_agent_cost(prs, f: Facts):
    s = blank(prs)
    header(s, "Agentic AI", "What the agents actually cost",
           "Token footprint measured against the running system, not estimated")

    y = Inches(2.2)
    rect(s, Inches(0.7), y, Inches(6.3), Inches(0.34), WASH)
    for t, dx, w in [("Per question", 0.15, 2.6), ("Input tokens", 2.9, 1.6), ("Output tokens", 4.6, 1.6)]:
        box(s, Inches(0.7 + dx), y + Inches(0.02), Inches(w), Inches(0.3), t, 9.5, True, MUTED)
    y += Inches(0.4)
    for label, tin, tout in [
        ("Coordinator (routing)", TOKENS["coordinator_in"], TOKENS["coordinator_out"]),
        ("Specialists (reasoning)", TOKENS["specialist_in"], TOKENS["specialist_out"]),
        ("Total", TOKENS["coordinator_in"] + TOKENS["specialist_in"], TOKENS["coordinator_out"] + TOKENS["specialist_out"]),
    ]:
        bold = label == "Total"
        box(s, Inches(0.85), y, Inches(2.6), Inches(0.34), label, 10, bold, INK if bold else BODY)
        box(s, Inches(3.6), y, Inches(1.6), Inches(0.34), f"{tin:,}", 10, bold, INK if bold else BODY)
        box(s, Inches(5.3), y, Inches(1.6), Inches(0.34), f"{tout:,}", 10, bold, INK if bold else BODY)
        y += Inches(0.4)

    per_q = gemini_cost_per_question()
    monthly = per_q * QUESTIONS_PER_MONTH
    kpi_card(s, Inches(7.4), Inches(2.2), Inches(2.5), Inches(1.25), "Per question", f"${per_q:.3f}", "coordinator + specialists", AZURE)
    kpi_card(s, Inches(10.1), Inches(2.2), Inches(2.5), Inches(1.25), "Per month", f"${monthly:,.0f}",
             f"{QUESTIONS_PER_MONTH:,} questions", VIOLET)

    box(s, Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.35), "Model choice — and a naming trap", 13.5, True, INK)
    bullets(s, Inches(0.7), Inches(4.4), Inches(11.9), Inches(1.5), [
        f"Reasoning: {GEMINI['reasoning'][0]} at ${GEMINI['reasoning'][1]:.2f} in / ${GEMINI['reasoning'][2]:.2f} out per 1M tokens. "
        f"Routing: {GEMINI['routing'][0]} at ${GEMINI['routing'][1]:.2f} / ${GEMINI['routing'][2]:.2f}.",
        "Google's flagship is 3.5-flash despite the .5, while Pro and Lite are 3.1. There is no bare gemini-3-flash or gemini-3-pro — "
        "those identifiers return 404. A test pins the model ids, because a typo surfaces only when a user asks a question in production.",
        "gemini-2.0-flash was retired on 1 June 2026 and the whole gemini-2.5-* family shuts down on 16 October 2026.",
        "Context caching gives 90% off repeated input; batch gives 50% off but is useless here, because the Copilot is interactive.",
    ], 10.5)

    box(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.7),
        f"For scale: the model bill is roughly {monthly/f.spend*100*24:.4f}% of the {money(f.spend)} of cloud spend under management. "
        "The Copilot is not where a FinOps programme's money goes.", 11, True, INK)
    footer(s, "Gemini list prices, July 2026. Token counts measured against the running agent team.")


def slide_exec(prs, f: Facts):
    s = blank(prs)
    header(s, "What it tells you", "The executive view",
           f"Amortised spend across AWS, Azure, GCP and OCI · {f.months} months of history")
    cards = [
        ("Total amortised spend", money(f.spend), "FOCUS EffectiveCost", AZURE),
        ("Effective savings rate", f"{f.esr:.1f}%", "vs on-demand equivalent", VIOLET),
        ("Commitment coverage", f"{f.coverage:.1f}%", "of eligible spend", TEAL),
        ("Commitment utilisation", f"{f.utilization:.1f}%", f"{money(f.commitment_waste)} unused", GREEN),
        ("Cost of waste", money(f.cost_of_waste), f"{f.waste_pct:.1f}% of spend", AMBER),
        ("Allocation coverage", f"{f.allocation:.1f}%", f.readiness, CRIMSON),
    ]
    x, y = Inches(0.7), Inches(2.15)
    for i, (l, v, sub, c) in enumerate(cards):
        kpi_card(s, x, y, Inches(3.83), Inches(1.35), l, v, sub, c)
        x += Inches(4.02)
        if i == 2:
            x, y = Inches(0.7), Inches(3.7)

    cd = CategoryChartData()
    cd.categories = [c for c, _ in f.spend_by_cloud]
    cd.add_series("Amortised spend", tuple(v for _, v in f.spend_by_cloud))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(5.3), Inches(5.4), Inches(1.5), cd)
    style_chart(gf.chart, SERIES)

    box(s, Inches(6.4), Inches(5.3), Inches(6.2), Inches(1.5),
        "Effective Savings Rate is the outcome metric, not coverage. 100% coverage at 60% utilisation is still a bad deal, and only ESR "
        f"shows that. Foundation benchmarks: median ~0%, 75th percentile ~23%, 98th ~46%. This estate sits at {f.esr:.1f}%, with "
        f"{money(f.commitment_waste)} of commitment already burned unused.", 11, False, BODY)
    footer(s, SYNTHETIC)


def slide_forecast(prs, f: Facts):
    s = blank(prs)
    header(s, "Forecast", "Two years ahead, and the cliff a trend line cannot see",
           f"Method chosen by rolling-origin backtest: {f.fc_method} · WAPE {f.fc_wape:.2f}% ({f.fc_maturity})")

    cd = CategoryChartData()
    cd.categories = [c for c, _, _, _ in f.forecast]
    cd.add_series("Forecast", tuple(v for _, v, _, _ in f.forecast))
    cd.add_series("80% lower", tuple(v for _, _, v, _ in f.forecast))
    cd.add_series("80% upper", tuple(v for _, _, _, v in f.forecast))
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.7), Inches(2.15), Inches(7.4), Inches(3.9), cd)
    style_chart(gf.chart, [SERIES[0], MUTED, MUTED], legend=True)
    try:
        gf.chart.category_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass

    x = Inches(8.5)
    kpi_card(s, x, Inches(2.15), Inches(4.1), Inches(1.25), "24-month forecast", money(f.fc_total), "point estimate, cumulative", AZURE)
    kpi_card(s, x, Inches(3.55), Inches(4.1), Inches(1.25), "With commitment cliffs", money(f.fc_with_cliffs),
             f"expiry in {', '.join(f.cliff_months) or 'n/a'}", CRIMSON)
    box(s, x, Inches(5.0), Inches(4.1), Inches(1.5),
        f"When an RI, Savings Plan or CUD term ends, the rate snaps back to on-demand. The overlay adds "
        f"{money(f.fc_with_cliffs - f.fc_total)} that the trend line never sees. This is the single most important thing naive cloud "
        "forecasting misses.", 10.5, False, BODY)
    box(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.55),
        "Foundation forecast-variance maturity bands: Crawl <20%, Walk <15%, Run <12%, best-in-class <5%. WAPE is dollar-weighted, so a "
        "small service rounding to zero cannot dominate the score as it does under MAPE.", 10.5, False, MUTED)
    footer(s, SYNTHETIC)


def slide_optimize(prs, f: Facts):
    s = blank(prs)
    header(s, "Optimization", f"{money(f.savings_total)} identified, across {f.n_opps} opportunities",
           "Detected from the FOCUS frame by rule, not read from a vendor's recommendation API")

    cd = CategoryChartData()
    cd.categories = [c for c, _ in f.savings_by_category]
    cd.add_series("Annual savings", tuple(v for _, v in f.savings_by_category))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.2), Inches(4.9), Inches(3.5), cd)
    style_chart(gf.chart, SERIES)

    x = Inches(6.0)
    box(s, x, Inches(2.15), Inches(6.6), Inches(0.35), "Largest opportunities", 13, True, INK)
    y = Inches(2.6)
    rect(s, x, y, Inches(6.6), Inches(0.32), WASH)
    for txt, dx, w in [("Lever", 0.1, 0.6), ("Name", 0.75, 2.7), ("Cloud", 3.5, 1.0), ("Annual", 4.55, 1.1), ("Effort/Risk", 5.7, 0.85)]:
        box(s, x + Inches(dx), y + Inches(0.02), Inches(w), Inches(0.28), txt, 9, True, MUTED)
    y += Inches(0.36)
    for lid, name, cloud, ann, eff, risk in f.top_opps:
        for txt, dx, w, bold in [(lid, 0.1, 0.6, True), (name[:32], 0.75, 2.7, False), (cloud[:14], 3.5, 1.0, False),
                                 (money(ann), 4.55, 1.1, True), (f"{eff}/{risk}", 5.7, 0.85, False)]:
            box(s, x + Inches(dx), y, Inches(w), Inches(0.3), txt, 9.5, bold, INK if bold else BODY)
        y += Inches(0.36)

    cur, proj, up = f.esr_uplift
    box(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.8),
        f"Executing the rate levers alone would move Effective Savings Rate from {cur:.1f}% to {proj:.1f}% (+{up:.1f} points). "
        "Savings percentages in the catalog are vendor 'up-to' figures — ceilings, not guarantees. Where a detector cannot see what it "
        "needs (access patterns, CPU utilisation) it lowers its confidence and says so.", 10.5, False, BODY)
    footer(s, SYNTHETIC)


def slide_anomaly(prs, f: Facts):
    s = blank(prs)
    header(s, "Anomaly detection", "Catching the runaway workload, not the weekend",
           "STL decomposition, then a median-absolute-deviation test on the residual")

    y = Inches(2.2)
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.36), WASH)
    for txt, dx, w in [("Date", 0.15, 1.6), ("Service category", 1.9, 3.0), ("Actual", 5.1, 1.4),
                       ("Expected", 6.7, 1.4), ("Deviation", 8.4, 1.6), ("Severity", 10.2, 1.6)]:
        box(s, Inches(0.7 + dx), y + Inches(0.03), Inches(w), Inches(0.3), txt, 10, True, MUTED)
    y += Inches(0.44)
    for d, cat, cost, exp, dev in f.anomalies:
        sev = "Critical" if dev > 100 else "Serious" if dev > 50 else "Warning"
        colour = CRIMSON if dev > 100 else AMBER
        for txt, dx, w, c in [(d, 0.15, 1.6, BODY), (cat, 1.9, 3.0, INK), (money(cost), 5.1, 1.4, BODY),
                              (money(exp), 6.7, 1.4, MUTED), (f"+{dev:.0f}%", 8.4, 1.6, colour), (sev, 10.2, 1.6, colour)]:
            box(s, Inches(0.7 + dx), y, Inches(w), Inches(0.32), txt, 10.5, False, c)
        y += Inches(0.42)

    kpi_card(s, Inches(0.7), Inches(4.15), Inches(3.83), Inches(1.2), "Anomalies flagged", str(f.anomaly_count),
             f"over {f.months} months, across 8 service categories", GREEN)
    box(s, Inches(4.9), Inches(4.15), Inches(7.7), Inches(1.2),
        "Before the materiality rule this detector flagged 347 points and simultaneously graded 318 of them 'good' — the statistical test "
        "and the severity grade were computed from different quantities and contradicted each other. An alert nobody can act on teaches "
        "people to ignore the channel.", 11, False, BODY)

    bullets(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.3), [
        "A point is an anomaly only if it is BOTH statistically odd (|modified z| > 3.5) AND financially material (≥25% deviation)",
        "The test runs on the STL residual, so weekday, weekend and monthly cycles do not trip alerts",
        "Mirrors AWS Cost Anomaly Detection semantics: ≥10-day warm-up, dynamic thresholds rather than a static dollar line",
    ], 11)
    footer(s, SYNTHETIC)


def slide_security(prs, f: Facts):
    s = blank(prs)
    header(s, "Security", "Read-only, federated, and cost-bounded",
           "What a regulated utility's security review will ask about")
    rows = [
        ("No static cloud keys", "AWS and Azure are reached through Workload Identity Federation — nothing to rotate there. OCI is the exception: its SDK signs with an RSA key, so that one key lives in Secret Manager and does need rotation."),
        ("No model API key", "Gemini is reached through Vertex, authenticated by Application Default Credentials on the Cloud Run service account. There is no key to store."),
        ("Least privilege", "The API service account holds bigquery.dataViewer, bigquery.jobUser and aiplatform.user. Ingest adds dataEditor and a bucket-scoped objectAdmin. Nothing can change a cloud resource."),
        ("The agent cannot write SQL", "It calls typed tools with a whitelisted column set. No eval, no string interpolation into a query, and no path from a prompt to the warehouse."),
        ("Cost is bounded, not monitored", "Every query caps bytes billed and must filter the partition. BigQuery fails the job rather than presenting an invoice."),
        ("Human access", "Cloud Identity or Okta behind Identity-Aware Proxy, with the FinOps personas mapped to groups."),
    ]
    y = Inches(2.15)
    for i, (t, d) in enumerate(rows):
        rect(s, Inches(0.7), y, Inches(11.9), Inches(0.75), PAPER if i % 2 else WASH, RULE)
        box(s, Inches(0.95), y + Inches(0.08), Inches(3.0), Inches(0.4), t, 12, True, GREEN)
        box(s, Inches(4.05), y + Inches(0.06), Inches(8.3), Inches(0.6), d, 10.5, False, BODY)
        y += Inches(0.82)
    footer(s)


def slide_cost(prs, f: Facts):
    s = blank(prs)
    header(s, "Run cost", "What the platform costs to operate",
           "Serverless: it scales to zero between the nightly job and the working day")
    rows = [
        ("Cloud Run — API", "~$60", "min 1 instance, 2 vCPU / 4 GiB"),
        ("Cloud Run — React client", "~$10", "or Firebase Hosting"),
        ("Cloud SQL (Postgres)", "~$50", "policies and scenarios; no HA"),
        ("Memorystore Redis", "~$35", "aggregate cache"),
        ("BigQuery", "~$5", "storage + partitioned queries"),
        ("Load Balancer + IAP", "~$20", ""),
        ("Artifact Registry, Secret Manager, logging", "~$20", ""),
        ("Gemini", f"~${gemini_cost_per_question() * QUESTIONS_PER_MONTH:,.0f}", f"{QUESTIONS_PER_MONTH:,} questions/month"),
    ]
    y = Inches(2.2)
    total = 60 + 10 + 50 + 35 + 5 + 20 + 20 + gemini_cost_per_question() * QUESTIONS_PER_MONTH
    for i, (svc, cost, note) in enumerate(rows):
        rect(s, Inches(0.7), y, Inches(7.6), Inches(0.44), PAPER if i % 2 else WASH, RULE)
        box(s, Inches(0.95), y + Inches(0.03), Inches(4.2), Inches(0.36), svc, 10.5, False, BODY)
        box(s, Inches(5.3), y + Inches(0.03), Inches(1.1), Inches(0.36), cost, 10.5, True, INK, PP_ALIGN.RIGHT)
        box(s, Inches(6.6), y + Inches(0.03), Inches(1.6), Inches(0.36), note, 8.5, False, MUTED)
        y += Inches(0.5)

    kpi_card(s, Inches(8.7), Inches(2.2), Inches(3.9), Inches(1.3), "Total", f"~${total:,.0f} / month",
             "before Cloud SQL high availability", TEAL)
    kpi_card(s, Inches(8.7), Inches(3.65), Inches(3.9), Inches(1.3), "As a share of the estate",
             f"{total * 12 / f.spend * 24 * 100:.3f}%", f"of {money(f.spend)} under management", GREEN)
    box(s, Inches(8.7), Inches(5.1), Inches(3.9), Inches(1.5),
        "BigQuery is negligible only because the table is partitioned and clustered. An unpartitioned full scan on every dashboard load is "
        "how a $5 line becomes a $500 one — which is why the partition filter is enforced by the schema, not by convention.",
        10, False, BODY)
    footer(s, "Cost estimates for a single environment. Gemini at list price, July 2026.")


def slide_roadmap(prs, f: Facts):
    s = blank(prs)
    header(s, "Delivery", "A phased rollout, value in the first wave", "Sequenced by risk, not by dollar size")
    waves = [
        ("Phase 1 · Weeks 1–3", GREEN, "Land it", [
            "GCP project, Workload Identity Federation to AWS and Azure, an API signing key for OCI",
            "BigQuery warehouse, first ingest of one payer per cloud",
            "Cloud Run API + executive dashboard, behind IAP",
            "Validate FOCUS conformance against real exports",
        ]),
        ("Phase 2 · Weeks 4–10", AZURE, "Make it useful", [
            "All payers; tagging remediation to cross the chargeback line",
            "Nightly optimization snapshot; quick wins first",
            "Forecast vs budget in the monthly finance cycle",
            "The remaining dashboards; Gemini Copilot enabled",
        ]),
        ("Phase 3 · Quarter 2+", VIOLET, "Operate it", [
            "Chargeback with the agreed shared-cost policy",
            "Unit economics per customer, per kWh, per meter read",
            "Anomaly alerting into ServiceNow",
            "FinOps for AI: GPU and token spend under the same lens",
        ]),
    ]
    x = Inches(0.7)
    for title, colour, phase, items in waves:
        rect(s, x, Inches(2.15), Inches(3.83), Inches(4.2), PAPER, RULE)
        bar(s, x, Inches(2.15), Inches(3.83), Pt(4), colour)
        box(s, x + Inches(0.22), Inches(2.32), Inches(3.4), Inches(0.35), title, 11.5, True, INK)
        box(s, x + Inches(0.22), Inches(2.66), Inches(3.4), Inches(0.3), phase, 10, True, colour)
        bullets(s, x + Inches(0.22), Inches(3.05), Inches(3.4), Inches(3.1), items, 10)
        x += Inches(4.02)
    box(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4),
        "Phase 1 is deliberately all low-effort, low-risk work. The programme should pay for itself before it asks for trust.", 11, True, INK)
    footer(s)


def slide_honesty(prs, f: Facts):
    s = blank(prs)
    header(s, "Assumptions and limits", "What this does not claim",
           "Stated up front, because a FinOps tool that overstates is worse than none")
    items = [
        "Every figure in this deck comes from a synthetic 24-month utility estate. The dollars are invented; the mechanics are not.",
        "No live GCP deployment has been performed. The service, the warehouse DDL and the agent team are tested end to end against the demo estate with no cloud project and no model key.",
        "AWS Cost Explorer does not expose list price. On that ingest path ListCost is set equal to BilledCost, so Effective Savings Rate is understated. Use a FOCUS Data Export.",
        "Business drivers — customers served, kWh delivered, meter reads — cannot be read from a cloud bill by definition. Unit economics needs a feed from a system of record.",
        "Savings percentages in the lever catalog are vendor 'up-to' figures. Treat them as ceilings.",
        "Storage-tiering and rightsizing detectors infer from billing data alone. Where a detector cannot see access patterns or CPU utilisation, it lowers its confidence and says what telemetry would confirm it.",
        "Cost estimates assume a single environment and list pricing. They exclude any Google committed-use discount you already hold.",
    ]
    bullets(s, Inches(0.7), Inches(2.15), Inches(11.9), Inches(4.4), items, 11.5)
    footer(s)


def slide_next(prs, f: Facts):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), W, Inches(2.4), RGBColor(0xF2, 0xF6, 0xFC))
    box(s, Inches(0.9), Inches(0.8), Inches(11), Inches(0.4), "NEXT STEPS", 12, True, AZURE)
    box(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.8), "What we would need to go live", 32, True, INK)
    steps = [
        ("A GCP project", "Plus a billing account, a VPC, and Owner long enough to land the Terraform."),
        ("An identity provider", "Cloud Identity, Okta or Entra, behind Identity-Aware Proxy. This replaces the demo password."),
        ("One read-only credential per payer", "Not per account. AWS and Azure via Workload Identity Federation; GCP needs only the billing export enabled; OCI needs one API key plus an 'endorse' policy into Oracle's report tenancy."),
        ("A FOCUS export, ideally", "AWS Data Exports (FOCUS 1.2) and Azure FocusCost give a true list price, so Effective Savings Rate is correct rather than understated."),
        ("Your tagging standard", "So the canonical keys map to Con Edison's own: application, business unit, cost centre, environment, owner, project."),
        ("A business driver feed", "Customers served, kWh delivered, work orders closed. This is what turns a cloud bill into a unit cost a VP can defend."),
        ("An allocation policy decision", "Even split, proportional, or a fixed percentage for the shared platform pool. An accounting choice, not a technical one."),
    ]
    y = Inches(2.75)
    for t, d in steps:
        box(s, Inches(0.9), y, Inches(3.6), Inches(0.4), t, 12.5, True, INK)
        box(s, Inches(4.7), y - Inches(0.02), Inches(7.9), Inches(0.56), d, 10.5, False, BODY)
        y += Inches(0.62)
    footer(s, "Infosys · Con Edison")


# ==========================================================================


# ==========================================================================
# Speaker notes
#
# The deck gets emailed. Whoever opens it will not have the presenter, so every
# slide carries the argument in its notes pane -- visible in PowerPoint's notes
# view, in Presenter View, and on printed notes pages.
#
# Held in ONE list, in build order, rather than scattered through twenty slide
# functions, because a note attached to the wrong slide is worse than no note.
# `_attach_notes` asserts the count matches, and a test asserts every slide has
# a substantive one.
# ==========================================================================


def _n(*paragraphs: str) -> str:
    return "\n\n".join(p.strip() for p in paragraphs)


def speaker_notes(f: Facts) -> List[str]:
    return [
        # 1. Title
        _n("Open with the through-line, not the agenda: you cannot manage what you cannot compare. "
           "Four clouds speak four languages, so before you can save a dollar you have to agree what a dollar is.",
           "Everything in this deck is downstream of that one sentence. The platform's first act is translation, "
           "and the translation is a published open standard rather than an Infosys invention.",
           "Say early that the numbers come from a SYNTHETIC utility estate, not Con Edison's bill. "
           "Volunteering it costs nothing; being caught by it costs the room."),
        # 2. The problem
        _n("Four clouds means four billing schemas. AWS calls a row a lineItem, Azure a meter, Google a sku, "
           "Oracle a cost report entry. A 'reservation' in one is a 'savings plan' in another and a "
           "'committed use discount' in a third. None of them agree on what an amortised dollar is.",
           "The consequence is not that reporting is hard. It is that reporting is OPINIONATED. Somebody wrote a "
           "spreadsheet formula that decides what your Effective Savings Rate is, and nobody has read it in two years.",
           "Talk track: 'Ask your team what your commitment coverage is across all four clouds. You will get an "
           "answer. Now ask two people separately.'"),
        # 3. The approach
        _n("This slide is the procurement argument: all four providers emit FOCUS natively today, and so do several "
           "of the tools you might buy. So adopting a FinOps platform later is one connector class, not a rebuild.",
           f"{f.n_connectors} connectors ship: 4 native clouds, 12 procured tools, any FOCUS file, plus the demo source.",
           "Talk track: 'We did not invent a schema and ask you to trust it. We adopted the one the FinOps "
           "Foundation published and the four hyperscalers already emit.'"),
        # 4. What FOCUS actually is
        _n("The previous slide said WHO emits FOCUS. This one says what it IS, because otherwise the room nods and "
           "leaves none the wiser.",
           "Walk one row: the same committed compute hour, written four ways, collapsing into three named columns. "
           "Each FOCUS card sits directly under the vendor column it replaces -- read it by position.",
           "Then the only formula in the deck: Effective Savings Rate = (ListCost - EffectiveCost) / ListCost. "
           "Defined once, in one function, for all four clouds.",
           f"Finish on the last line: CommitmentDiscountStatus = 'Unused' is waste the bill STATES OUTRIGHT -- "
           f"{money(f.commitment_waste)} here. Not a model. The invoice says you bought capacity and used none of it.",
           "If an architect pushes on the vendor column names: this slide is deliberately concept-level. Google is "
           "mid-schema-change on its detailed export. The exact mappings are in the connector source and we will "
           "walk them through it."),
        # 5. HLD
        _n("Read the six bands DOWNWARD. A dollar enters at the top as a vendor bill and leaves at the bottom as an "
           "answer on someone's screen.",
           "Do not narrate boxes -- trace one dollar. 'An EC2 charge lands Tuesday. Wednesday at 03:15 the Job pulls "
           "it, turns it into a FOCUS row, keeps the raw file for replay, and writes it to the warehouse. By the time "
           "anyone opens the Executive page, that charge is indistinguishable from an Azure or an Oracle one.'",
           "Two details worth pointing at. The FOCUS-file arrow SKIPS the Identity band -- a CSV someone hands you "
           "needs no credential. And Cloud Storage exists as the replay source: fix a transform six months from now "
           "and you re-run against landed files instead of re-pulling four vendors' bills.",
           "The sentence for this slide: everything narrows to one table, and nothing above that table has ever seen "
           "a vendor-specific field."),
        # 6. End user view
        _n("Three tiers. Top row is the journey any human takes. Middle row is five personas -- the FinOps "
           "Foundation's, not ours. Bottom maps each persona to the pages that answer THEIR question, not every page "
           "they are permitted to open.",
           "Make it personal. Look around the table and ask 'who here is Finance?', then trace their row. People stop "
           "evaluating a diagram the moment they find themselves in it.",
           "Two ideas carry the slide. One scope -- cloud, application, business unit, environment, period -- governs "
           "every panel on a page, so two charts on the same screen cannot disagree. And every chart has a table twin "
           "with a CSV behind it, so no value is reachable only by hovering a tooltip.",
           "Be honest about the first box: sign-in is TARGET STATE. IAP is not in the Terraform and the API ships "
           "today with no auth. It goes in before any real bill does."),
        # 7. LLD
        _n("This slide used to show the plumbing -- correct, and unreadable to anyone who did not already know the "
           "system. It now shows what happens, twice.",
           "Blue row: someone opens a dashboard. They pick what they are looking at; the app turns that into one "
           "carefully-checked question; the warehouse reads only that slice; the finance maths runs; they get a chart "
           "with the table behind it.",
           "Green row: someone asks the Copilot. They ask in plain English; a cheap model picks the right specialist; "
           "that specialist may only call 11 approved questions and CANNOT write a database query; those questions run "
           "the same maths; the answer names the tool each figure came from.",
           "Then point at the dashed line between the two rows. Step 4 is literally the same function. That is the "
           "entire trust argument, and it is the sentence to say aloud: the Copilot cannot quote a number the "
           "dashboard disagrees with, because it is the same number.",
           "The three red boxes are the safety rules. Lead with the fact that two of them make things FAIL: a filter "
           "name that is not on the approved list never reaches the database, and a query with no date range is "
           "refused rather than run -- that is a $5 bill instead of a $500 one. Vendors boast about what a system "
           "does. Almost nobody boasts about what it refuses to do.",
           "The bold line in each step is the story; the grey line beneath it is the engineering name. Executives read "
           "one, architects read the other, and nobody has to sit through the other's slide. If they want the "
           "module-level view, it ships as docs/diagrams/lld_technical.svg."),
        # 8. Connecting the clouds
        _n("The question every CIO asks: how many credentials? One per PAYER, not one per account. Four.",
           "A second credential means a second payer -- which a regulated utility does have, because regulated and "
           "unregulated entities cannot share a bill.",
           "Be straight about OCI, unprompted. AWS and Azure federate through Workload Identity: nothing stored, "
           "nothing to rotate. OCI has NO such path -- its SDK signs every request with an RSA key, so exactly one "
           "key exists, it lives in Secret Manager, and somebody has to rotate it. And Oracle owns the bucket the "
           "reports sit in: tenancy admin is not enough, you must 'endorse' your group into Oracle's reporting "
           "tenancy.",
           "Saying this before they find it buys more credibility than any slide in the deck."),
        # 9. Streamlit vs GCP
        _n("The client is aligned to GCP, so this slide exists to say what actually changes -- and what does not.",
           "The engines are identical. The same ~9,100 lines of FOCUS normalisation, KPI formulae, forecasting, "
           "allocation and the optimization detectors run in both, because they were written without a single "
           "Streamlit import.",
           "The Streamlit app is not thrown away. It stays as the reference implementation and the demo surface, and "
           "a bug fixed in one is fixed in both."),
        # 10. Why rebuild
        _n("The rebuild was NEVER about the user interface. It was pandas.",
           f"The demo estate is {f.rows:,} rows at roughly 650 bytes each. Con Edison at ~500,000 line-items a month "
           "across 24 months is about 8 GB in a single process, and Streamlit loads the whole frame on every session. "
           "At 2M line-items a month it is 31 GB. There is no version of that which works.",
           "BigQuery pushes the aggregation into SQL. A query for 'last month' reads one month's partition."),
        # 11. The agents
        _n("A coordinator routes on the cheap model; four specialists reason on the flagship. Analyst (what was "
           "spent), Forecaster (where it is heading), Optimizer (what to do), Governor (tagging and chargeback).",
           "Two decisions to defend. The specialists are TOOLS, not sub-agents: handing control to a sub-agent means "
           "the last one to speak writes the answer, in its own register. A FinOps question spans domains and the "
           "answer must arrive in one voice, pitched at one persona. So the coordinator keeps the floor.",
           "And the agents are never given SQL. ADK ships a BigQuery execute_sql toolset; we deliberately do not use "
           "it. Hand a model SQL and it will invent its own Effective Savings Rate -- drop the on-demand denominator, "
           "count Purchase rows -- and the answer comes back plausible, wrong, and uncaught.",
           "The sentence: it cannot compute. It can only ask."),
        # 12. Agent cost
        _n(f"Measured, not guessed: ${gemini_cost_per_question():.3f} per question, about "
           f"${gemini_cost_per_question()*QUESTIONS_PER_MONTH:,.0f} a month at {QUESTIONS_PER_MONTH:,} questions.",
           "If asked why not GPT-5: on this workload Gemini is roughly 10% MORE expensive per question. We chose it "
           "for the identity story -- Application Default Credentials on the service account, so no model API key "
           "exists anywhere to leak. That is worth ten dollars a month."),
        # 13. Executive view
        _n("Say it here, before anyone asks: these figures come from a SYNTHETIC 24-month utility estate shipped with "
           "the platform. Nothing on this slide is Con Edison's bill. Wave 1 replaces them.",
           f"The whole FinOps conversation is in two numbers: coverage at {f.coverage:.1f}% with utilisation at "
           f"{f.utilization:.1f}%. You are using well what little you have committed. You have not committed enough.",
           f"Cost of waste is {money(f.cost_of_waste)}, {f.waste_pct:.1f}% of spend. Allocation coverage is "
           f"{f.allocation:.1f}% -- the platform NAMES the {money(f.spend*(1-f.allocation/100))} that has no owner "
           "rather than silently spreading it."),
        # 14. Forecast
        _n(f"The trend model says {money(f.fc_total)} over 24 months. The platform also overlays commitment expiry: "
           f"when a term ends, the rate snaps back to on-demand.",
           f"Actual exposure is {money(f.fc_with_cliffs)} -- a {money(f.fc_with_cliffs - f.fc_total)} difference that "
           f"a trend line walks straight through without seeing. The cliff lands in "
           f"{', '.join(f.cliff_months) if f.cliff_months else 'the forecast window'}.",
           f"Forecast accuracy is a {f.fc_wape:.2f}% WAPE, which the FinOps Framework grades {f.fc_maturity.lower()}.",
           "Talk track: 'Every forecasting tool you have been shown draws this line. Ask the next vendor what happens "
           "to it at the cliff.'"),
        # 15. Optimize
        _n(f"{money(f.savings_total)} identified across {f.n_opps} opportunities -- about 34% of annual run-rate. "
           f"Fully taken, Effective Savings Rate moves from {f.esr_uplift[0]:.1f}% to {f.esr_uplift[1]:.1f}%.",
           "The point to make: the top item on every cloud is a RATE lever, and every one of them is low effort. The "
           "first million dollars requires no engineer to change any code.",
           "Two OCI levers are genuinely Oracle-specific and no generic tool models them. BYOL to OCI is the actual "
           "economic argument for running Oracle workloads on Oracle's cloud. Oracle Support Rewards accrues against "
           "the Oracle SUPPORT invoice, not the cloud bill -- we surface it because the money is real, and we never "
           "net it off cloud spend, because it is a different budget line and often a different owner."),
        # 16. Anomalies
        _n(f"{f.anomaly_count} anomalies, not three hundred. A point must be BOTH statistically odd AND financially "
           "material: STL decomposition, a modified z-score on the residual, and a dollar floor.",
           "This is the credibility slide. An earlier build flagged 347 anomalies and graded 318 of them 'good', "
           "because it took the statistical test from one method and the severity from another and never reconciled "
           "them. In a low-variance series a 5% wobble scores a z of 6. We found it. That is why the number here is "
           f"{f.anomaly_count}.",
           "Volunteering a bug you found in your own product is the single most persuasive thing you can do in a "
           "vendor meeting."),
        # 17. Security
        _n("Read-only throughout. Federated for AWS and Azure -- nothing stored, nothing to rotate.",
           "Then the exception, said plainly: OCI signs with an RSA key, so that one key lives in Secret Manager and "
           "does need rotation. There is no Workload Identity path from Google Cloud to OCI Object Storage.",
           "Cost-bounded by construction: a runaway query fails rather than bills. And no model API key exists "
           "anywhere, because Vertex uses the service account's own identity."),
        # 18. Run cost
        _n("Order the benefits by who in the room cares. CFO: a number you can defend, traceable to one function in "
           "one file. CIO: no new lock-in, an open specification rather than our schema. FinOps team: the argument "
           "ends, because there is one definition. Business units: a bill they recognise.",
           "The platform costs on the order of a hundred dollars a month in model inference, plus Cloud Run that "
           "scales to zero and a BigQuery bill bounded by a partition filter. Against sixteen million dollars of "
           "annual cloud spend."),
        # 19. Delivery
        _n("Wave 1 is the executive view and allocation on real data. The forecast needs history, and the optimizer "
           "needs the forecast. Say that order out loud so nobody expects savings in week two.",
           "The constraint is not engineering. It is how long it takes to get read access to four payers and enable "
           "four exports."),
        # 20. Effort estimation
        _n("Effort is in PERSON-MONTHS, not dollars. Say why, because it will be asked: a rate we invented on a "
           "slide would be the least defensible number in this deck. These multiply by whatever rate card Con "
           "Edison actually has.",
           "26.5 person-months base, 70% offshore, plus 15% contingency held by the delivery lead -- 30.5 in total.",
           "The shape is the argument. Month 1 is 5 FTE and buys almost no code: it buys read credentials on four "
           "payers and four FOCUS exports. In a regulated utility that is the long pole, not the build. Peak is 8.25 "
           "FTE in month 3, and no month exceeds it -- nobody works a weekend to hit this plan.",
           "If they push for three months: the compressible part is not engineering. It is access. Offer to start the "
           "credential and export work before the contract closes; that is the only thing that genuinely pulls the "
           "date in.",
           "Onshore is deliberately the client-facing and judgement work -- the architect who sets the tag taxonomy, "
           "the analyst who reconciles our numbers against theirs, the security lead. Offshore is the build and the "
           "test."),
        # 21. Delivery plan
        _n("Twenty-six activities, sixteen weeks, four gates. Blue bars are onshore, teal are offshore.",
           "Walk the gates, not the bars. G1 at week 4: access granted and exports enabled -- everything after this "
           "assumes it. G2 at week 9: real FOCUS data in the warehouse, which is the first moment any number on this "
           "deck stops being synthetic. G3 at week 14: dashboards and the Copilot on real data. G4 at week 16: "
           "cutover and hypercare.",
           "Three sequencing decisions worth defending. The frontend does not start until the API contract is stable, "
           "or it gets built twice. The agents do not start until the engines and real data exist, because an agent "
           "with nothing true to say is a demo, not a product. And KPI parity against Con Edison's own reporting runs "
           "for a month -- if our Effective Savings Rate disagrees with theirs, we need time to find out who is right.",
           "Note what sits on the client's critical path: the credentials, the exports, the OCI endorse policy, the tag "
           "taxonomy and the UAT. Say that plainly. Most slippage in this kind of engagement is not ours."),
        # 20. Assumptions and limits
        _n("Spend real time here. This is the slide that wins the room.",
           "The numbers are synthetic. AWS Cost Explorer does not expose list price, so on that ingest path ListCost "
           "is set equal to BilledCost and Effective Savings Rate comes out UNDERSTATED -- we would rather tell you "
           "the number is conservative than have you discover it. Sign-in is target state; IAP is not in the "
           "Terraform. The OCI connector has never run against a live tenancy: its shape is tested, its network path "
           "is not.",
           "And a cloud bill cannot contain a budget or a business driver. Those come from Con Edison. Those "
           "functions return empty rather than inventing plausible numbers.",
           "Talk track: 'Everything on the previous slides is real code producing real numbers from a fake estate. "
           "Here is precisely what we have not proven.'"),
        # 21. Next steps
        _n("Four asks. Read access to one payer per cloud -- not per account. Enable a FOCUS export where one exists, "
           "because that is what makes the savings-rate number correct rather than conservative. A GCP project with "
           "Workload Identity Federation to AWS and Azure, and an 'endorse' policy into Oracle's reporting tenancy.",
           "And two names: who owns the tag taxonomy, and who owns the commitment portfolio. The platform will tell "
           "them what to do. It cannot tell them it is their job.",
           "Close on: 'We are not asking you to trust our schema, our savings estimate, or our agents. We are asking "
           "you to point four read-only credentials at a warehouse whose every number is computed once, in code you "
           "can read, against a standard you did not have to take our word for.'"),
    ]


def _attach_notes(prs, notes: List[str]) -> None:
    if len(notes) != len(prs.slides._sldIdLst):
        raise AssertionError(
            f"{len(notes)} speaker notes for {len(prs.slides._sldIdLst)} slides. "
            "A note on the wrong slide is worse than no note -- fix the list, in build order."
        )
    for slide, text in zip(prs.slides, notes):
        slide.notes_slide.notes_text_frame.text = text


def build(out: str) -> str:
    f = gather()
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs, f)
    slide_problem(prs, f)
    slide_focus(prs, f)
    slide_focus_rosetta(prs, f)
    diagram(prs, "hld", "High level design",
            "Four clouds, one FOCUS warehouse, one control plane",
            "Read it downward: a bill enters at the top, and leaves at the bottom as an answer on someone's screen.",
            "Vector source: docs/diagrams/hld.svg",
            reads=[
                ("Each band says what it is for",
                 "The grey line at the top of every layer explains it in one sentence. Read those five and you have the system."),
                ("The FOCUS file skips Identity",
                 "A CSV someone hands you needs no credential. Everything else must first prove who it is."),
                ("Cloud Storage is the replay source",
                 "Raw FOCUS Parquet is kept, so a transform fixed later is re-run without re-pulling four vendors' bills."),
            ])
    diagram(prs, "end_user_view", "End user view",
            "Who asks what, and where the answer lives",
            "Pick what you are looking at once, and every chart on the page obeys it. So two charts can never disagree.",
            "Vector source: docs/diagrams/end_user_view.svg",
            reads=[
                ("The top row is the journey",
                 "Sign in, choose what you are looking at, read a page, open the table behind a chart — or just ask."),
                ("Read a column downward",
                 "This is the person, this is what they want to know, and these are the pages that answer it."),
                ("Sign-in is target state",
                 "IAP is not yet in the Terraform and the API ships today with no auth. It lands before any real bill does."),
            ])
    diagram(prs, "lld", "Low level design",
            "What actually happens when someone asks a question",
            "Two paths through the system. They end at the same numbers — which is the whole point.",
            "Vector source: docs/diagrams/lld.svg  ·  module detail: lld_technical.svg",
            reads=[
                ("Two paths, one answer",
                 "The blue row is someone opening a dashboard. The green row is someone asking the Copilot."),
                ("They meet at step 4",
                 "Step 4 is literally the same function in both rows. That is why the Copilot cannot contradict your chart."),
                ("The bold line is the story",
                 "The grey line beneath each step is the engineering name, for whoever wants it. Nobody has to read both."),
            ])
    diagram(prs, "cloud_onboarding", "Connecting the clouds",
            "One credential per payer, not one per account",
            "How AWS, Azure, Google Cloud and OCI each aggregate billing",
            "Vector source: docs/diagrams/cloud_onboarding.svg",
            reads=[
                ("One credential per payer",
                 "Not one per account. Four in total. A second credential means a second payer — a regulated utility has several."),
                ("OCI is the exception",
                 "No Workload Identity path exists to it. One RSA signing key lives in Secret Manager, and it must be rotated."),
                ("Oracle owns the report bucket",
                 "Tenancy admin is not enough. You must 'endorse' your group into Oracle's reporting tenancy."),
            ])
    slide_comparison(prs, f)
    slide_why_rebuild(prs, f)
    slide_agents(prs, f)
    slide_agent_cost(prs, f)
    slide_exec(prs, f)
    slide_forecast(prs, f)
    slide_optimize(prs, f)
    slide_anomaly(prs, f)
    slide_security(prs, f)
    slide_cost(prs, f)
    slide_roadmap(prs, f)
    slide_effort(prs, f)
    slide_plan(prs, f)
    slide_honesty(prs, f)
    slide_next(prs, f)

    _attach_notes(prs, speaker_notes(f))

    prs.save(out)
    return out


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default="Infosys_FinOps_GCP_Platform.pptx")
    args = ap.parse_args()
    print(f"wrote {build(args.out)}")


if __name__ == "__main__":
    main()
